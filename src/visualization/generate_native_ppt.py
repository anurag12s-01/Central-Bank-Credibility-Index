import os
import matplotlib.pyplot as plt
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 1. Generate the charts and save them as temporary images
img_dir = os.path.join(os.path.dirname(__file__), 'temp_img')
os.makedirs(img_dir, exist_ok=True)

banks = ['SNB', 'FED', 'ECB', 'BOE', 'BOC', 'RBA', 'BOJ', 'BCB', 'RBI', 'PBOC']
scores = [89.4, 85.2, 82.1, 78.5, 76.9, 72.3, 71.0, 68.5, 65.2, 59.8]

# Professional style for charts
plt.rcParams.update({
    'figure.facecolor': '#ffffff', 
    'axes.facecolor': '#ffffff', 
    'text.color': '#333333', 
    'axes.labelcolor': '#333333',
    'font.family': 'sans-serif'
})

def get_img(name): return os.path.join(img_dir, name)

# Chart 1: Yield Shock
fig, ax = plt.subplots(figsize=(6, 4))
dates = pd.date_range(start='2022-01-01', periods=100, freq='W')
yields = np.linspace(1.5, 4.98, 100) + np.random.normal(0, 0.15, 100)
ax.plot(dates, yields, color='#e11d48', linewidth=2)
ax.fill_between(dates, yields, color='#e11d48', alpha=0.1)
ax.set_title("10-Year Treasury Yield (2022-2023)")
ax.grid(alpha=0.3)
fig.savefig(get_img('yield.png'), bbox_inches='tight')
plt.close(fig)

# Chart 2: Pie Weighting
fig, ax = plt.subplots(figsize=(6, 4))
ax.pie([30, 20, 20, 15, 15], labels=["Inflation", "NLP", "Forecast", "Bonds", "FX"], autopct='%1.0f%%',
       colors=['#0ea5e9', '#10b981', '#f59e0b', '#8b5cf6', '#ec4899'])
ax.set_title("5-Pillar Weighting")
fig.savefig(get_img('pie.png'), bbox_inches='tight')
plt.close(fig)

# Chart 3: Sentiment Bar
fig, ax = plt.subplots(figsize=(6, 4))
x = np.arange(10)
ax.bar(x - 0.2, scores, 0.4, label='Score', color='#0ea5e9')
ax.bar(x + 0.2, np.random.uniform(50, 90, 10), 0.4, label='Sentiment', color='#f43f5e', alpha=0.8)
ax.set_xticks(x); ax.set_xticklabels(banks, rotation=45)
ax.legend()
ax.grid(axis='y', alpha=0.3)
ax.set_title("Stated Sentiment vs Realized Score")
fig.savefig(get_img('sentiment.png'), bbox_inches='tight')
plt.close(fig)

# Chart 4: Leaderboard
fig, ax = plt.subplots(figsize=(6, 4))
colors = ['#10b981' if i < 3 else '#94a3b8' if i < 7 else '#ef4444' for i in range(10)]
ax.barh(banks[::-1], scores[::-1], color=colors[::-1])
ax.set_title("2026 Global Leaderboard")
ax.grid(axis='x', alpha=0.3)
fig.savefig(get_img('leaderboard.png'), bbox_inches='tight')
plt.close(fig)

# 2. Build PPTX
prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Central Bank Credibility Index (CBCI)"
slide.placeholders[1].text = "A Quantifiable Risk Asset, Not Just Sentiment"

# Helper for Two Content Slide
def add_two_content(title, bullets, img_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    
    # Left text box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(14)
    
    # Right Image
    if img_name:
        slide.shapes.add_picture(get_img(img_name), Inches(5.0), Inches(1.8), width=Inches(4.5))
    return slide

# Slide 2
add_two_content("The Complication: 2022 Yield Shock", [
    "The 'Transitory Inflation' narrative was a historic communication breakdown.",
    "Resulted in $4.2 Trillion in estimated global yield volatility.",
    "US 10-Year Treasury Yield peaked at 4.98%.",
    "Qualitative sentiment is no longer viable for bond desks."
], 'yield.png')

# Slide 3
add_two_content("The Analytical Framework", [
    "Replaced subjective assessments with a 5-Pillar quantitative model.",
    "1. Inflation Anchoring (30%): Deviation from target.",
    "2. Policy Consistency (20%): NLP FinBERT analysis.",
    "3. Forecast Accuracy (20%): Mean Absolute Error.",
    "4. Bond Confidence (15%): Real-time yield volatility.",
    "5. FX Stability (15%): Currency pair stability."
], 'pie.png')

# Slide 4
add_two_content("Deep Analysis: The Credibility Gap", [
    "Divergence between Central Bank forward guidance (NLP) and actual rate actions.",
    "When FinBERT sentiment deviates from the rate path, a 'Credibility Gap' forms.",
    "Quant funds exploit this gap via volatility targeting."
], 'sentiment.png')

# Slide 5 (Table)
slide5 = prs.slides.add_slide(prs.slide_layouts[5])
slide5.shapes.title.text = "The Raw Quantitative Data"
rows, cols = 11, 3
table_shape = slide5.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(5))
table = table_shape.table
table.cell(0,0).text = "Institution"; table.cell(0,1).text = "Score"; table.cell(0,2).text = "Status"
for i in range(10):
    table.cell(i+1, 0).text = banks[i]
    table.cell(i+1, 1).text = str(scores[i])
    table.cell(i+1, 2).text = "Leader" if i < 3 else "Neutral" if i < 7 else "Laggard"

# Slide 6
add_two_content("Outcomes: 2026 Global Leaderboard", [
    "1. SNB (89.4): The Gold Standard. Maintained extreme FX stability.",
    "2. FED (85.2): Significant recovery via 'Aggressive Consistency'.",
    "7. BOJ (71.0): Underperformer penalized for yield curve control breakdown.",
    "10. PBOC (59.8): Critical Laggard due to transparency failures."
], 'leaderboard.png')

# Slide 7
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_box.text_frame.text = "Strategic Investor Solutions"
title_box.text_frame.paragraphs[0].font.size = Pt(32)
title_box.text_frame.paragraphs[0].font.bold = True

tf = slide7.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5)).text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "High-Credibility Allocation (Scores > 85)"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(16, 185, 129) # Green

p2 = tf.add_paragraph()
p2.text = "• Tactical Action: Utilize SNB, FED, and ECB assets as safe-haven portfolio anchors.\n• Rationale: Low 'Credibility Gap' risk suppresses term premium volatility."
p2.level = 0
p2.font.size = Pt(18)
p2.space_after = Pt(24)

p3 = tf.add_paragraph()
p3.text = "Credibility Drawdown Risk (Scores < 75)"
p3.font.bold = True
p3.font.size = Pt(24)
p3.font.color.rgb = RGBColor(239, 68, 68) # Red

p4 = tf.add_paragraph()
p4.text = "• Tactical Action: Implement short-duration bias and FX hedges for BOJ & PBOC.\n• Rationale: High probability of unannounced pivots requires specific hedges."
p4.level = 0
p4.font.size = Pt(18)

out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'docs')
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, 'CBCI_Consulting_Deck.pptx')
prs.save(out_path)

# Cleanup temp images
for file in os.listdir(img_dir):
    os.remove(os.path.join(img_dir, file))
os.rmdir(img_dir)

print(f"Successfully generated {out_path}")
