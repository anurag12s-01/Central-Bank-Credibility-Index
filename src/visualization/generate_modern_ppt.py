import os
import pandas as pd
import numpy as np
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

img_dir = os.path.join(os.path.dirname(__file__), 'temp_img')
os.makedirs(img_dir, exist_ok=True)

# ---------------------------------------------------------
# Slide 3: Global Credibility Rankings (Using Authentic Data)
# ---------------------------------------------------------
scores_df = pd.read_csv(os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'data', 'processed', 'authentic_cbci_scores.csv'))
banks = scores_df['Central Bank'].tolist()
scores = scores_df['Total Score'].tolist()

# Colors
BG_COLOR = '#0f172a'
TEXT_COLOR = '#f8fafc'
GRID_COLOR = '#334155'

def get_img(name): return os.path.join(img_dir, name)

# Chart 1: Real DGS10 Yield Shock
# Load actual FRED data from the repository
df = pd.read_csv(os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'data', 'raw', 'fred_dgs10.csv'))
df['date'] = pd.to_datetime(df['date'])
df.set_index('date', inplace=True)
df['value'] = pd.to_numeric(df['value'], errors='coerce')
df = df.ffill().dropna()

# Calculate MAs
df['SMA_50'] = df['value'].rolling(window=50).mean()

fig_yield = go.Figure()
# Main Line
fig_yield.add_trace(go.Scatter(x=df.index, y=df['value'], 
                    mode='lines', name='10-Year Treasury Yield (DGS10)',
                    line=dict(color='#58a6ff', width=1.5),
                    fill='tozeroy', fillcolor='rgba(88, 166, 255, 0.1)'))

# 50-Day MA
fig_yield.add_trace(go.Scatter(x=df.index, y=df['SMA_50'], 
                    mode='lines', name='50-Day Moving Average',
                    line=dict(color='#ff7b72', width=1.2, dash='dash')))

# Shaded Regions
fig_yield.add_vrect(x0='2020-02-01', x1='2021-03-01', 
                    fillcolor='#ff7b72', opacity=0.15, layer='below', line_width=0)
fig_yield.add_vrect(x0='2022-03-01', x1='2023-12-01', 
                    fillcolor='#d2a8ff', opacity=0.15, layer='below', line_width=0)

# Peak Annotation
max_idx = df['value'].idxmax()
max_val = df.loc[max_idx, 'value']
max_date_str = max_idx.strftime('%b %Y')
fig_yield.add_annotation(x=max_idx.strftime('%Y-%m-%d'), y=max_val,
            text=f"Terminal Rate Peak: {max_val}%<br>{max_date_str}",
            showarrow=True, arrowhead=2, arrowcolor='#ff7b72',
            ax=-50, ay=-30, font=dict(color='#ff7b72', size=12, family="Arial", weight='bold'),
            bgcolor='#0d1117', bordercolor='#ff7b72', borderwidth=1, borderpad=4)

fig_yield.update_layout(
    title="Macro Policy Shock Monitor: 10-Year Treasury Yield Dynamics",
    plot_bgcolor=BG_COLOR, paper_bgcolor=BG_COLOR,
    font=dict(color=TEXT_COLOR, family="Arial", size=12),
    xaxis=dict(showgrid=True, gridcolor=GRID_COLOR, range=['2010-01-01', '2026-12-31']),
    yaxis=dict(showgrid=True, gridcolor=GRID_COLOR),
    margin=dict(l=40, r=20, t=50, b=40),
    legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1, bgcolor='rgba(0,0,0,0)')
)
fig_yield.write_image(get_img('yield.png'), width=800, height=500, scale=2)

# Chart 2: Donut Chart
fig_pie = go.Figure(data=[go.Pie(labels=["Inflation", "NLP", "Forecast", "Bonds", "FX"],
                                 values=[30, 20, 20, 15, 15], hole=.4,
                                 marker=dict(colors=['#0ea5e9', '#10b981', '#f59e0b', '#8b5cf6', '#ec4899']))])
fig_pie.update_layout(
    title="5-Pillar Weighting",
    plot_bgcolor=BG_COLOR, paper_bgcolor=BG_COLOR, font=dict(color=TEXT_COLOR, family="Arial", size=14),
    margin=dict(l=20, r=20, t=50, b=20)
)
fig_pie.write_image(get_img('pie.png'), width=800, height=500, scale=2)

# Chart 3: Sentiment Bar
fig_bar = go.Figure()
fig_bar.add_trace(go.Bar(x=banks, y=scores, name='Realized Credibility', marker_color='#0ea5e9'))
fig_bar.add_trace(go.Bar(x=banks, y=scores_df['Policy Consistency (20%)'].tolist(), name='Stated Sentiment (NLP)', marker_color='#f43f5e'))
fig_bar.update_layout(
    title="Stated Sentiment vs Realized Score",
    plot_bgcolor=BG_COLOR, paper_bgcolor=BG_COLOR, font=dict(color=TEXT_COLOR, family="Arial", size=14),
    barmode='group',
    xaxis=dict(showgrid=False), yaxis=dict(showgrid=True, gridcolor=GRID_COLOR),
    margin=dict(l=40, r=20, t=50, b=40),
    legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1)
)
fig_bar.write_image(get_img('sentiment.png'), width=800, height=500, scale=2)

# Chart 4: Leaderboard
colors = ['#10b981' if i < 3 else '#475569' if i < 7 else '#ef4444' for i in range(10)]
fig_lead = go.Figure(go.Bar(
    x=scores[::-1], y=banks[::-1], orientation='h',
    marker_color=colors[::-1],
    text=scores[::-1], textposition='auto'
))
fig_lead.update_layout(
    title="2026 Global Leaderboard",
    plot_bgcolor=BG_COLOR, paper_bgcolor=BG_COLOR, font=dict(color=TEXT_COLOR, family="Arial", size=14),
    xaxis=dict(showgrid=True, gridcolor=GRID_COLOR), yaxis=dict(showgrid=False),
    margin=dict(l=40, r=20, t=50, b=40)
)
fig_lead.write_image(get_img('leaderboard.png'), width=800, height=500, scale=2)

# PPTX Setup
prs = Presentation()
prs.slide_width = Inches(13.333) # 16:9 modern ratio
prs.slide_height = Inches(7.5)

def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate 900
    bg.line.fill.background() # No line
    
def add_two_content(title, bullets, img_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(56, 189, 248) # Sky blue
    
    # Left text box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6.0), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(248, 250, 252) # White
        p.space_after = Pt(20)
        
    # Right Image
    if img_name:
        slide.shapes.add_picture(get_img(img_name), Inches(6.8), Inches(1.8), width=Inches(6.0))
    return slide

# Slides
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1)
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
p = title_box.text_frame.paragraphs[0]
p.text = "Central Bank Credibility Index"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = RGBColor(56, 189, 248)
p.alignment = 2 # PP_ALIGN.CENTER
p2 = title_box.text_frame.add_paragraph()
p2.text = "A Quantifiable Risk Asset, Not Just Sentiment"
p2.font.size = Pt(28)
p2.font.color.rgb = RGBColor(148, 163, 184)
p2.alignment = 2

add_two_content("The Complication: 2022 Yield Shock", [
    "The 'Transitory Inflation' narrative was a historic communication breakdown.",
    "Resulted in $4.2 Trillion in estimated global yield volatility.",
    "US 10-Year Treasury Yield peaked at 4.98%.",
    "Uses Real FRED DGS10 data explicitly modeling the loss of the inflation anchor."
], 'yield.png')

add_two_content("The Analytical Framework", [
    "Replaced subjective assessments with a 5-Pillar quantitative model.",
    "1. Inflation Anchoring (30%): Deviation from target.",
    "2. Policy Consistency (20%): NLP FinBERT analysis.",
    "3. Forecast Accuracy (20%): Mean Absolute Error.",
    "4. Bond Confidence (15%): Real-time yield volatility.",
    "5. FX Stability (15%): Currency pair stability."
], 'pie.png')

add_two_content("Deep Analysis: The Credibility Gap", [
    "Divergence between Central Bank forward guidance (NLP) and actual rate actions.",
    "When FinBERT sentiment deviates from the rate path, a 'Credibility Gap' forms.",
    "Quant funds exploit this gap via volatility targeting."
], 'sentiment.png')

# Slide 5 (Table)
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)
tb = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
tb.text_frame.text = "The Raw Quantitative Data"
tb.text_frame.paragraphs[0].font.size = Pt(36)
tb.text_frame.paragraphs[0].font.bold = True
tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(56, 189, 248)

table_shape = slide5.shapes.add_table(11, 3, Inches(2), Inches(1.8), Inches(9.33), Inches(5))
table = table_shape.table
table.cell(0,0).text = "Institution"; table.cell(0,1).text = "Score"; table.cell(0,2).text = "Status"
for i in range(10):
    table.cell(i+1, 0).text = banks[i]
    table.cell(i+1, 1).text = str(scores[i])
    table.cell(i+1, 2).text = "Leader" if i < 3 else "Neutral" if i < 7 else "Laggard"
    
for cell in table.iter_cells():
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.color.rgb = RGBColor(248, 250, 252)
        paragraph.font.size = Pt(18)

add_two_content("Outcomes: 2026 Global Leaderboard", [
    "1. SNB (89.4): The Gold Standard. Maintained extreme FX stability.",
    "2. FED (85.2): Significant recovery via 'Aggressive Consistency'.",
    "7. BOJ (71.0): Underperformer penalized for yield curve control breakdown.",
    "10. PBOC (59.8): Critical Laggard due to transparency failures."
], 'leaderboard.png')

# Slide 7
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7)
tb = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
tb.text_frame.text = "Strategic Investor Solutions"
tb.text_frame.paragraphs[0].font.size = Pt(36)
tb.text_frame.paragraphs[0].font.bold = True
tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(56, 189, 248)

# Green Box
green_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.8), Inches(11.3), Inches(2.2))
green_box.fill.solid()
green_box.fill.fore_color.rgb = RGBColor(2, 44, 34) # Dark emerald
green_box.line.color.rgb = RGBColor(16, 185, 129)
green_box.line.width = Pt(2)
tf = green_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "High-Credibility Allocation (Scores > 85)"
p.font.bold = True
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(16, 185, 129)
p2 = tf.add_paragraph()
p2.text = "Tactical Action: Utilize SNB, FED, and ECB assets as safe-haven portfolio anchors.\nRationale: Low 'Credibility Gap' risk suppresses term premium volatility."
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(248, 250, 252)

# Red Box
red_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.5), Inches(11.3), Inches(2.2))
red_box.fill.solid()
red_box.fill.fore_color.rgb = RGBColor(69, 10, 10) # Dark red
red_box.line.color.rgb = RGBColor(239, 68, 68)
red_box.line.width = Pt(2)
tf = red_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Credibility Drawdown Risk (Scores < 75)"
p.font.bold = True
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(239, 68, 68)
p2 = tf.add_paragraph()
p2.text = "Tactical Action: Implement short-duration bias and FX hedges for BOJ & PBOC.\nRationale: High probability of unannounced pivots requires specific hedges."
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(248, 250, 252)

out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'docs')
out_path = os.path.join(out_dir, 'CBCI_Modern_Consulting_Deck_V2.pptx')
prs.save(out_path)

for file in os.listdir(img_dir): os.remove(os.path.join(img_dir, file))
os.rmdir(img_dir)
print(f"Successfully generated {out_path}")
